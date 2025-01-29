from .utils import count_levels
from pptx import Presentation
from pptx import util
from datetime import date, datetime
import json


class Prs:
    """Expects factory_json (json str) with the following structure:
        {
        'Methodology': {
            'title': 'Methodology',
            'text': 'Paragraph'
        },

        'Key Findings': {
            'title': 'Key Findings',
            'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
        },

        'Publisher Performance': {
            'Final Bias Rating': {
                'title': 'Slide Title',
                'chart_title': 'Chart Title',
                'chart_filepath': 'prs-demo-files/bar_chart_bias_rating.png',
                'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
            }, 
            'Bias Category': {
                'title': 'Slide Title',
                'chart_title': 'Chart Title',
                'chart_filepath': 'prs-demo-files/bar_chart_bias_category.png',
                'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
            },
            'Topics vs. Final Bias Rating': {
                'title': 'Slide Title',
                'chart_title': 'Chart Title',
                'chart_filepath': 'prs-demo-files/heatmap_topic_bias_rating.png',
                'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
            },
        },

        'Publisher Comparison': {
            'Tendency to Commit Bias': {
                'title': 'Slide Title',
                'chart_title': 'Chart Title',
                'chart_filepath': 'prs-demo-files/odds_bias_rating.png',
                'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
            },
        },

        'Use Cases': {
            'Very Biased': {
                'title': 'Article Title',
                'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
            }
        },

        'Conclusion': {
            'title': 'Conclusion',
            'bullets': 'Bullet 1 \n Bullet 2 \n Bullet 3'
        }
    }
    """
    def __init__(self, pptx_template_filepath, factory_json=None, placeholder_json=None):
        # Initialize PPT template
        self.prs = Presentation(pptx_template_filepath)

        # Convert png + LLM json results into a dict
        factory_dict = None
        if factory_json is not None:
            factory_dict = json.loads(factory_json)

        # Convert placeholder text contents into a dict
        placeholder_dict = None
        if placeholder_json is not None:
            placeholder_dict = json.loads(placeholder_json)

        # Create content_dict
        self.content_dict = {}
        for d in [factory_dict, placeholder_dict]:
            if d is not None:
                self.content_dict.update(d)

        # Slide type indicator --
        # This will change whenever slide master of PPTX template is changed
        self.layout_dict = {
            'Title': 0,
            'Section': 1,
            'Title, Text': 2,
            'Title, Content': 3,
            'Title, Text, Content, Picture': 4,
            'Title, Content, Picture': 5,
            'Title, Picture, Picture': 6,
            'Blank': 7
        }

    class AdjustedPicture:
        """Fits picture into the container size in the PPT slide"""

        def __init__(self, pptx_placeholder_object):
            self.pptx_placeholder_object = pptx_placeholder_object
            self.adjusted_pptx_placeholder_object = None

        def fit_to_container(self):
            picture = self.pptx_placeholder_object
            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0

            available_width = picture.width
            available_height = picture.height
            image_width, image_height = picture.image.size
            placeholder_aspect_ratio = float(available_width) / float(available_height)
            image_aspect_ratio = float(image_width) / float(image_height)

            # Get initial image placeholder left and top positions
            pos_left, pos_top = picture.left, picture.top

            picture.crop_top = 0
            picture.crop_left = 0
            picture.crop_bottom = 0
            picture.crop_right = 0

            # If the placeholder is "wider" in aspect, shrink the picture width while
            # Maintaining the image aspect ratio
            if placeholder_aspect_ratio > image_aspect_ratio:
                picture.width = int(image_aspect_ratio * available_height)
                picture.height = available_height

            # Otherwise shrink the height
            else:
                picture.height = int(available_width/image_aspect_ratio)
                picture.width = available_width

            # Set the picture left and top position to the initial placeholder one
            picture.left, picture.top = pos_left, pos_top

            self.adjusted_pptx_placeholder_object = picture

            return self.adjusted_pptx_placeholder_object


    def _render_title_slide(self, title='Briefing Pack', subtitle=date.today().strftime("%B %d, %Y")):
        slide_layout = self.prs.slide_layouts[self.layout_dict['Title']]
        slide = self.prs.slides.add_slide(slide_layout)

        placeholder = slide.placeholders[0]
        placeholder.text = title
        placeholder = slide.placeholders[1]
        placeholder.text = subtitle

        return self.prs


    def _render_section_slide(self, title='Section Name', subtitle=''):
        slide_layout = self.prs.slide_layouts[self.layout_dict['Section']]
        slide = self.prs.slides.add_slide(slide_layout)

        placeholder = slide.placeholders[0]
        placeholder.text = title
        placeholder = slide.placeholders[1]
        placeholder.text = subtitle

        return self.prs


    def _render_paragraph_slide(self, title='', text=''):
        slide_layout = self.prs.slide_layouts[self.layout_dict['Title, Text']]
        slide = self.prs.slides.add_slide(slide_layout)

        placeholder = slide.placeholders[0]
        placeholder.text= title
        placeholder = slide.placeholders[1]
        placeholder.text= text

        return self.prs


    def _render_bullets_slide(self, title='', bullets=''):
        slide_layout = self.prs.slide_layouts[self.layout_dict['Title, Content']]
        slide = self.prs.slides.add_slide(slide_layout)

        placeholder = slide.placeholders[0]
        placeholder.text= title
        placeholder = slide.placeholders[4]
        placeholder.text= bullets

        return self.prs


    def _render_chartbullets_slide(self, title='', chart_title='', chart_filepath='', bullets=''):
        slide_layout = self.prs.slide_layouts[self.layout_dict['Title, Text, Content, Picture']]
        slide = self.prs.slides.add_slide(slide_layout)

        placeholder = slide.placeholders[0]
        placeholder.text = title
        placeholder = slide.placeholders[1]
        placeholder.text = chart_title
        placeholder = slide.placeholders[4]
        placeholder.text = bullets
        placeholder = slide.placeholders[13]
        
        if chart_filepath != '':
            picture = placeholder.insert_picture(chart_filepath)
            picture_adj = self.AdjustedPicture(picture).fit_to_container()

        return self.prs


    def _render_imagebullets_slide(self, title='', image_filepath='', bullets=''):
        slide_layout = self.prs.slide_layouts[self.layout_dict['Title, Content, Picture']]
        slide = self.prs.slides.add_slide(slide_layout)

        placeholder = slide.placeholders[0]
        placeholder.text = title
        placeholder = slide.placeholders[2]
        placeholder.text = bullets
        placeholder = slide.placeholders[13]
        if image_filepath != '':
            picture = placeholder.insert_picture(image_filepath)
            picture_adj = self.AdjustedPicture(picture).fit_to_container()

        self.prs = self.prs

        return self.prs


    def add_Title_section(self, title='', subtitle=[], use_json=False):
        if use_json:
            key = 'Title'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                subtitle = d['subtitle']
                self.prs = self._render_title_slide(title, subtitle)
            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    subtitle = v['subtitle']
                    self.prs = self._render_title_slide(title, subtitle)
        else:
            start_date, end_date = subtitle
            dt_start = datetime.strptime(start_date, '%Y-%m-%d')
            dt_end = datetime.strptime(end_date, '%Y-%m-%d')
            if dt_start.year == dt_end.year and dt_start.month == dt_end.month:
                subtitle = dt_start.strftime('%B %Y')
            else:
                subtitle = f'{dt_start.strftime('%B %Y')} to {dt_end.strftime('%B %Y')}'
            self.prs = self._render_title_slide(title, subtitle)

        return self.prs

    
    def add_Section_section(self, title='', subtitle='', use_json=False):
        if use_json:
            key = 'Section'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                subtitle = d['subtitle']
                self.prs = self._render_section_slide(title, subtitle)
            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    subtitle = v['subtitle']
                    self.prs = self._render_section_slide(title, subtitle)
        else:
            self.prs = self._render_section_slide(title, subtitle)

        return self.prs


    def add_Introduction_section(self, title='', text='', use_json=False):
        if use_json:
            key = 'Introduction'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                text = d['text']
                self.prs = self._render_paragraph_slide(title, text)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    text = v['text']
                    self.prs = self._render_paragraph_slide(title, text)       
        else:
            self.prs = self._render_paragraph_slide(title, text)

        return self.prs


    def add_Methodology_section(self, title='', text='', use_json=False):
        if use_json:
            key = 'Methodology'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                text = d['text']
                self.prs = self._render_paragraph_slide(title, text)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    text = v['text']
                    self.prs = self._render_paragraph_slide(title, text)       
        else:
            self.prs = self._render_paragraph_slide(title, text)

        return self.prs

    
    def add_KeyFindings_section(self, title='', bullets='', use_json=False):
        if use_json:
            key = 'Key Findings'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                bullets = d['text']
                self.prs = self._render_bullets_slide(title, bullets)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    bullets = v['text']
                    self.prs = self._render_bullets_slide(title, bullets)       
        else:
            self.prs = self._render_bullets_slide(title, bullets)

        return self.prs


    def add_PublisherPerformance_section(self, title='', chart_title='', chart_filepath='', bullets='', use_json=False):
        if use_json:
            key = 'Publisher Performance Overview'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                chart_title = d['chart_title']
                chart_filepath = d['chart_filepath']
                bullets = d['bullets']
                self.prs = self._render_chartbullets_slide(title, chart_title, chart_filepath, bullets)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    chart_title = v['chart_title']
                    chart_filepath = v['chart_filepath']
                    bullets = v['bullets']
                    self.prs = self._render_chartbullets_slide(title, chart_title, chart_filepath, bullets)       
        else:
            self.prs = self._render_chartbullets_slide(title, chart_title, chart_filepath, bullets)

        return self.prs


    def add_PublisherComparison_section(self, title='', chart_title='', chart_filepath='', bullets='', use_json=False):
        if use_json:
            key = 'Publisher Comparison'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                chart_title = d['chart_title']
                chart_filepath = d['chart_filepath']
                bullets = d['bullets']
                self.prs = self._render_chartbullets_slide(title, chart_title, chart_filepath, bullets)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    chart_title = v['chart_title']
                    chart_filepath = v['chart_filepath']
                    bullets = v['bullets']
                    self.prs = self._render_chartbullets_slide(title, chart_title, chart_filepath, bullets)       
        else:
            self.prs = self._render_chartbullets_slide(title, chart_title, chart_filepath, bullets)

        return self.prs


    def add_UseCases_section(self, title='', image_filepath='', bullets='', use_json=False):
        if use_json:
            key = 'Case Studies'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                bullets = d['bullets']
                self.prs = self._render_imagebullets_slide(title, image_filepath, bullets)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    bullets = v['bullets']
                    self.prs = self._render_imagebullets_slide(title, image_filepath, bullets)       
        else:
            self.prs = self._render_imagebullets_slide(title, image_filepath, bullets)

        return self.prs


    def add_Conclusion_section(self, title='', bullets='', use_json=False):
        if use_json:
            key = 'Conclusions'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                bullets = d['text']
                self.prs = self._render_bullets_slide(title, bullets)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    bullets = v['text']
                    self.prs = self._render_bullets_slide(title, bullets)       
        else:
            self.prs = self._render_bullets_slide(title, bullets)

        return self.prs


    def add_Recommendations_section(self, title='', bullets='', use_json=False):
        if use_json:
            key = 'Recommendations'
            d = self.content_dict[key]

            if count_levels(d) == 1:
                title = d['title']
                bullets = d['bullets']
                self.prs = self._render_bullets_slide(title, bullets)

            elif count_levels(d) > 1:
                for k, v in d.items():
                    title = v['title']
                    bullets = v['bullets']
                    self.prs = self._render_bullets_slide(title, bullets)       
        else:
            self.prs = self._render_bullets_slide(title, bullets)

        return self.prs


    def save(self, filepath):
        self.prs.save(filepath)